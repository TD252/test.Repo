from pptx import Presentation
from pptx.util import Pt

def create_resume_interactive(filename="resume_interactive.pptx"):
    # Ask user for details
    name = input("Enter your full name: ")
    email = input("Enter your email: ")
    phone = input("Enter your phone number: ")
    skills = input("Enter your skills (comma separated): ").split(",")
    experience = input("Enter your experience (comma separated): ").split(",")
    education = input("Enter your education details: ")
    projects = input("Enter your projects (comma separated): ").split(",")

    # Create PowerPoint
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Name
    title_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(40), width=Pt(500), height=Pt(50))
    title_tf = title_box.text_frame
    title_tf.text = name

    # Contact Info
    contact_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(90), width=Pt(500), height=Pt(40))
    contact_tf = contact_box.text_frame
    contact_tf.text = f"Email: {email} | Phone: {phone}"

    # Skills
    skills_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(140), width=Pt(500), height=Pt(80))
    skills_tf = skills_box.text_frame
    skills_tf.text = "Skills:\n" + "\n".join([s.strip() for s in skills])

    # Experience
    exp_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(230), width=Pt(500), height=Pt(100))
    exp_tf = exp_box.text_frame
    exp_tf.text = "Experience:\n" + "\n".join([e.strip() for e in experience])

    # Education
    edu_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(340), width=Pt(500), height=Pt(60))
    edu_tf = edu_box.text_frame
    edu_tf.text = "Education:\n" + education

    # Projects
    proj_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(420), width=Pt(500), height=Pt(80))
    proj_tf = proj_box.text_frame
    proj_tf.text = "Projects:\n" + "\n".join([p.strip() for p in projects])

    # Save
    prs.save(filename)
    print(f"Resume created and saved as {filename}")

# Run the function
create_resume_interactive()
